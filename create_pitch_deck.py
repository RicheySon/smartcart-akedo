#!/usr/bin/env python3
"""
Script to convert markdown pitch deck to PowerPoint
Requires: python-pptx library
Install: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "SmartCart"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Intelligent Grocery Shopping Agent\nPowered by AI & Blockchain - BNB Chain Testnet"
    sp = subtitle_frame.paragraphs[0]
    sp.font.size = Pt(24)
    sp.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "The Problem"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Grocery Shopping is Broken"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    bullets = [
        "🛒 Waste: $218B wasted annually on expired groceries",
        "💰 Overspending: No budget tracking or price comparison",
        "⏰ Time-consuming: Manual inventory management",
        "📊 No Insights: No forecasting or smart recommendations",
        "🏪 Vendor Confusion: Hard to compare prices across platforms"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Solution"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "SmartCart - Your AI-Powered Grocery Assistant"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    
    features = [
        "📦 Tracks inventory automatically with expiration alerts",
        "🤖 Predicts run-out dates using ML linear regression",
        "💵 Manages budgets and prevents overspending",
        "🛍️ Compares prices across Amazon & Walmart",
        "✅ Approves transactions with risk assessment",
        "📝 Maintains audit logs for compliance"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.level = 0
    
    # Add more slides as needed...
    
    prs.save('SmartCart_Pitch_Deck.pptx')
    print("✅ Pitch deck created: SmartCart_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_pitch_deck()








